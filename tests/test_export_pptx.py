# tests/test_export_pptx.py
import io
from pptx import Presentation
from pptx.dml.color import RGBColor
from components.export_pptx import generate_pptx, _hex

DATA = {
    "profielen": {
        "engineer": {"kleur": "#0072B8", "titel": "QA Engineer"},
    },
    "blokken": {
        "engineer": [
            {
                "sectie": "Cloud",
                "items": [
                    {"id": f"item-{i}", "naam": f"Cursus {i}", "icon": "☁️",
                     "desc": f"Beschrijving {i}", "duur": "1 dag", "kern": True}
                    for i in range(7)
                ],
            }
        ]
    },
}

PLAN_LEEG = {
    "profiel": "engineer",
    "geselecteerd": [],
    "statussen": {},
}

PLAN_EEN_ITEM = {
    "profiel": "engineer",
    "geselecteerd": ["item-0"],
    "statussen": {"item-0": "afgerond"},
}

PLAN_ZES_ITEMS = {
    "profiel": "engineer",
    "geselecteerd": [f"item-{i}" for i in range(6)],
    "statussen": {f"item-{i}": "gepland" for i in range(6)},
}

PLAN_ZEVEN_ITEMS = {
    "profiel": "engineer",
    "geselecteerd": [f"item-{i}" for i in range(7)],
    "statussen": {f"item-{i}": "gepland" for i in range(7)},
}


def _parse(result: bytes) -> Presentation:
    return Presentation(io.BytesIO(result))


# ── _hex ──────────────────────────────────────────────────────────────────────

def test_hex_converteert_kleur():
    kleur = _hex("#E5007D")
    assert isinstance(kleur, RGBColor)
    assert kleur[0] == 0xE5
    assert kleur[1] == 0x00
    assert kleur[2] == 0x7D


def test_hex_zonder_hekje():
    assert _hex("0072B8") == _hex("#0072B8")


# ── generate_pptx: output ─────────────────────────────────────────────────────

def test_geeft_bytes_terug():
    result = generate_pptx(DATA, PLAN_LEEG)
    assert isinstance(result, bytes)
    assert len(result) > 0


def test_output_is_geldig_pptx():
    result = generate_pptx(DATA, PLAN_LEEG)
    prs = _parse(result)
    assert prs is not None


# ── slides tellen ─────────────────────────────────────────────────────────────

def test_leeg_plan_geeft_alleen_cover():
    prs = _parse(generate_pptx(DATA, PLAN_LEEG))
    assert len(prs.slides) == 1


def test_een_item_geeft_cover_plus_een_content_slide():
    prs = _parse(generate_pptx(DATA, PLAN_EEN_ITEM))
    assert len(prs.slides) == 2


def test_zes_items_passen_op_een_content_slide():
    prs = _parse(generate_pptx(DATA, PLAN_ZES_ITEMS))
    assert len(prs.slides) == 2


def test_zeven_items_vragen_twee_content_slides():
    prs = _parse(generate_pptx(DATA, PLAN_ZEVEN_ITEMS))
    assert len(prs.slides) == 3


# ── cover-slide inhoud ────────────────────────────────────────────────────────

def _teksten(slide) -> list[str]:
    return [
        run.text
        for shape in slide.shapes
        if shape.has_text_frame
        for para in shape.text_frame.paragraphs
        for run in para.runs
        if run.text.strip()
    ]


def test_cover_bevat_profielnaam():
    prs = _parse(generate_pptx(DATA, PLAN_EEN_ITEM))
    teksten = _teksten(prs.slides[0])
    assert any("QA Engineer" in t for t in teksten)


def test_cover_bevat_aantal_geselecteerde_blokken():
    prs = _parse(generate_pptx(DATA, PLAN_ZES_ITEMS))
    teksten = _teksten(prs.slides[0])
    assert any("6" in t for t in teksten)


# ── onbekend profiel valt terug op defaults ───────────────────────────────────

def test_onbekend_profiel_gooit_geen_fout():
    plan = {**PLAN_EEN_ITEM, "profiel": "bestaat-niet"}
    result = generate_pptx(DATA, plan)
    assert isinstance(result, bytes)
