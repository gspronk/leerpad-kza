# components/export_pptx.py
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def _hex(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


DONKER    = _hex("0D1117")
SURFACE   = _hex("161B22")
ROZE      = _hex("E5007D")
BLAUW     = _hex("0072B8")
GROEN     = _hex("39D353")
AMBER     = _hex("F0A500")
TEKST     = _hex("CDD9E5")
MUTED     = _hex("8B949E")
WIT       = _hex("F0F6FF")


def generate_pptx(data: dict, plan: dict) -> bytes:
    """Genereert PPTX van het huidige plan. Geeft bytes terug voor st.download_button."""
    prs = Presentation()
    prs.slide_width  = Inches(13.3)
    prs.slide_height = Inches(7.5)

    profiel_key  = plan.get("profiel", "engineer")
    profiel_meta = data.get("profielen", {}).get(profiel_key, {})
    profiel_kleur = _hex(profiel_meta.get("kleur", "#0072B8"))

    geselecteerd = set(plan.get("geselecteerd", []))
    statussen    = plan.get("statussen", {})

    # Bouw lookup van item-ID naar item
    alle_items = {
        item["id"]: {**item, "_profiel": prof}
        for prof, secties in data.get("blokken", {}).items()
        for sectie in secties
        for item in sectie.get("items", [])
    }

    sel_items = [alle_items[iid] for iid in geselecteerd if iid in alle_items]

    # Slide 1: Cover
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DONKER
    bg.line.fill.background()

    accent = slide.shapes.add_shape(1, 0, 0, Inches(0.4), prs.slide_height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = profiel_kleur
    accent.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(9), Inches(1.2))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"KZA Leerpad — {profiel_meta.get('titel', profiel_key)}"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = WIT

    tx2 = slide.shapes.add_textbox(Inches(0.7), Inches(3.8), Inches(9), Inches(0.6))
    p2 = tx2.text_frame.paragraphs[0]
    r2 = p2.add_run()
    r2.text = f"{len(geselecteerd)} bouwblokken geselecteerd"
    r2.font.size = Pt(14)
    r2.font.color.rgb = MUTED

    # Slide 2+: Blokken per 6 per slide
    ITEMS_PER_SLIDE = 6
    for slide_idx in range(0, len(sel_items), ITEMS_PER_SLIDE):
        batch = sel_items[slide_idx: slide_idx + ITEMS_PER_SLIDE]
        slide = prs.slides.add_slide(blank_layout)

        bg2 = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg2.fill.solid()
        bg2.fill.fore_color.rgb = DONKER
        bg2.line.fill.background()

        t = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12), Inches(0.5))
        p = t.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = f"Mijn Plan — {profiel_meta.get('titel', '')}"
        r.font.size = Pt(11)
        r.font.color.rgb = MUTED

        cols_count = 3
        card_w = Inches(4.0)
        card_h = Inches(2.8)
        margin_x = Inches(0.45)
        margin_y = Inches(0.75)
        gap_x = Inches(0.1)
        gap_y = Inches(0.1)

        for i, item in enumerate(batch):
            col = i % cols_count
            row = i // cols_count
            x = margin_x + col * (card_w + gap_x)
            y = margin_y + row * (card_h + gap_y)

            card = slide.shapes.add_shape(1, x, y, card_w, card_h)
            card.fill.solid()
            card.fill.fore_color.rgb = SURFACE
            card.line.color.rgb = _hex("21262D")
            card.line.width = Pt(0.5)

            status = statussen.get(item["id"], "gepland")
            st_kleur = GROEN if status == "afgerond" else (AMBER if status == "bezig" else MUTED)
            top_bar = slide.shapes.add_shape(1, x, y, card_w, Inches(0.06))
            top_bar.fill.solid()
            top_bar.fill.fore_color.rgb = st_kleur
            top_bar.line.fill.background()

            nt = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.18), card_w - Inches(0.3), Inches(0.5))
            np_ = nt.text_frame.paragraphs[0]
            nr = np_.add_run()
            nr.text = f"{item.get('icon','')} {item['naam']}"
            nr.font.size = Pt(11)
            nr.font.bold = True
            nr.font.color.rgb = WIT

            dt = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.75), card_w - Inches(0.3), Inches(1.4))
            dt.text_frame.word_wrap = True
            dp = dt.text_frame.paragraphs[0]
            dr = dp.add_run()
            dr.text = item.get("desc", "")[:120]
            dr.font.size = Pt(8)
            dr.font.color.rgb = MUTED

            ft = slide.shapes.add_textbox(x + Inches(0.15), y + card_h - Inches(0.45), card_w - Inches(0.3), Inches(0.3))
            fp = ft.text_frame.paragraphs[0]
            fr = fp.add_run()
            status_label = {"afgerond": "✓ Afgerond", "bezig": "◉ Bezig"}.get(status, "○ Gepland")
            fr.text = f"{item.get('duur','')}  ·  {status_label}"
            fr.font.size = Pt(8)
            fr.font.color.rgb = st_kleur

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
